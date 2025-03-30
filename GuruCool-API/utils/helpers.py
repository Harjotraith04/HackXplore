#utils/helpers.py
import os
import re
import requests
import logging
import threading
import time
import PyPDF2

from pathlib import Path
from pptx import Presentation

from langchain.schema import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import DirectoryLoader, PyPDFLoader, TextLoader

from utils.firebase import *

download_lock = threading.Lock()

def check_or_download_files(DATA_PATH: str, subject: str, unit: str, lecture_id: str) -> list:
    """
    Check if local documents exist for the given lecture ID. If they exist, load and return the documents.
    If not, fetch the document URLs from Firestore, download the files, and load and return the documents.
    Parameters:
    - DATA_PATH (str): The path to the data directory.
    - subject (str): The subject of the lecture.
    - unit (str): The unit of the lecture.
    - lecture_id (str): The ID of the lecture.
    Returns:
    - documents (list): A list of loaded documents, or None if an error occurred.
    """
    lecture_dir = os.path.join(DATA_PATH, subject, unit, lecture_id)
    os.makedirs(lecture_dir, exist_ok=True)

    if os.listdir(lecture_dir): 
        logging.info(f"Documents already exist for Lecture ID: {lecture_id}. Loading existing documents...")
        try:
            documents = load_documents(lecture_dir)
            return documents
        except Exception as e:
            logging.error(f"Error loading local documents: {e}")
            return None

    logging.info(f"No local documents found. Fetching document URLs and downloading files...")
    document_urls = fetch_document_urls_from_firestore(subject, unit, lecture_id)

    if not document_urls:
        logging.error(f"No document URLs found for Lecture ID: {lecture_id}")
        return None 
    
    download_files(document_urls, lecture_dir)
    
    try:
        documents = load_documents(lecture_dir)
        return documents
    except Exception as e:
        logging.error(f"Error loading downloaded documents: {e}")
        return None
        
def chunk_content(content: str, chunk_size=1000) -> str:
    return [content[i:i+chunk_size] for i in range(0, len(content), chunk_size)]

def doc_chunking(documents: list) -> list:
    """
    Generate content chunks from a list of documents.
    Args:
        documents (list): A list of documents.
    Returns:
        list: A list of content chunks.
    Raises:
        ValueError: If no documents are loaded.
    """
    try:
        if not documents:
            raise ValueError("No documents loaded.")
        
        combined_content = "\n\n".join(doc.page_content for doc in documents)
        content_chunks = chunk_content(combined_content)
        
        return content_chunks

    except Exception as e:
        print(f"Error generating quiz: {e}")
        return {}

def download_files(document_urls: list, save_directory: str) -> None:
    """
    Downloads files from the provided document data and saves them in the specified directory.
    If the file already exists, it will skip downloading.

    Args:
        document_urls (list): A list of URLs (strings) or dictionaries with 'url' and 'name'.
        save_directory (str): Directory where files will be saved.
    """
    logging.info(f"Downloading files from provided URLs")
    os.makedirs(save_directory, exist_ok=True)

    for idx, doc in enumerate(document_urls):
        try:
            with download_lock:
                if isinstance(doc, dict):
                    file_url = doc.get('url')
                else:
                    file_url = doc 

                if not file_url:
                    logging.error("Invalid document data: URL is empty.")
                    continue

                file_name_match = re.search(r'/([^/]+\.[a-zA-Z0-9]+)(?:\?|$)', file_url)
                file_name = file_name_match.group(1) if file_name_match else f'document_{idx}.txt'
                file_path = os.path.join(save_directory, file_name)

                if os.path.exists(file_path):
                    logging.info(f"File {file_name} already exists, skipping download.")
                    continue

                logging.info(f"Downloading {file_url} to {file_name}")
                response = requests.get(file_url)
                response.raise_for_status()
                
                with open(file_path, 'wb') as file:
                    file.write(response.content)

                logging.info(f"Downloaded and saved file: {file_name}")
        except requests.exceptions.RequestException as e:
            logging.error(f"Error downloading file from {file_url}: {e}")


def load_pptx_files(directory: str) -> list:
    """
    Load PPTX files from a given directory and return a list of documents.

    Args:
        directory (str): The directory path where the PPTX files are located.

    Returns:
        list: A list of documents, each containing the page content and metadata of a PPTX file.

    """
    logging.info(f"Loading PPTX files from directory: {directory}")
    pptx_documents = []
    for pptx_file in Path(directory).glob("*.pptx"):
        logging.info(f"Processing PPTX file: {pptx_file}")
        prs = Presentation(pptx_file)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        pptx_documents.append(Document(page_content="\n".join(text_runs), metadata={"source": str(pptx_file)}))
    logging.info(f"Loaded {len(pptx_documents)} PPTX documents.")
    return pptx_documents

def load_documents(directory: str) -> list:
    """
    Load documents from a given directory.
    Args:
        directory (str): The directory path where the documents are located.
    Returns:
        list: A list of split documents.
    Raises:
        FileNotFoundError: If the specified directory does not exist.
        Exception: If there is an error loading a text file.
    """
    logging.info(f"Loading documents from directory: {directory}")
    start_time = time.time()

    dir_path = Path(directory)

    txt_documents = []
    pdf_documents = []
    
    for file in dir_path.glob("*.txt"):
        try:
            loader = TextLoader(str(file))
            txt_documents.extend(loader.load())
            logging.info(f"Loaded text document: {file}")
        except Exception as e:
            logging.error(f"Error loading text file {file}: {e}")

    pdf_loader = DirectoryLoader(str(dir_path), glob="*.pdf", loader_cls=PyPDFLoader)
    pdf_documents = pdf_loader.load()
    logging.info(f"Loaded {len(pdf_documents)} PDF documents.")
    
    pptx_documents = load_pptx_files(directory)

    documents = txt_documents + pdf_documents + pptx_documents
    logging.info(f"Total documents loaded: {len(documents)}")

    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
    split_documents = text_splitter.split_documents(documents)
    logging.info(f"Documents split into {len(split_documents)} chunks.")

    end_time = time.time()
    logging.info(f"Documents loaded and processed in {end_time - start_time:.2f} seconds")

    return split_documents

def extract_text_from_books(book_url: str, save_directory: str) -> str:
    """
    Downloads the book from the provided Firebase URL and extracts text, saving the file in the specified directory.
    
    Args:
        book_url (str): The URL of the book in Firebase storage.
        save_directory (str): The directory where the book will be saved.
    
    Returns:
        str: The full extracted text from the PDF.
    """
    os.makedirs(save_directory, exist_ok=True)
    
    book_path = download_file_from_url(book_url, save_directory)

    full_text = ""
    print(f"Extracting text from book: {book_path}")

    pdf_reader = PyPDF2.PdfReader(book_path)
    for page_num in range(len(pdf_reader.pages)):
        page_text = pdf_reader.pages[page_num].extract_text()
        full_text += page_text + "\n\n"

    print("Text extraction from books complete.")
    return full_text

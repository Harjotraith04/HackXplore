import os
import sys
import warnings
import logging
import time
import pickle
import json
import whisper
from pathlib import Path

import requests
from flask import Flask, abort, request, jsonify, session
from flask_cors import CORS
from dotenv import load_dotenv
import firebase_admin
from firebase_admin import credentials, storage as firebase_storage, firestore
from langchain_ollama import OllamaLLM
from langchain_community.document_loaders import DirectoryLoader, PyPDFLoader, TextLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document
from sentence_transformers import SentenceTransformer
import faiss
from pptx import Presentation

warnings.filterwarnings("ignore", category=FutureWarning, module="whisper")
warnings.filterwarnings("ignore", category=UserWarning, module="whisper")

# --------------------------------------------------------------- ENV INIT --------------------------------------------------------------- 

load_dotenv()
STORAGE_BUCKET = os.getenv("STORAGE_BUCKET")
CREDENTIALS_PATH = os.getenv("CREDENTIALS_PATH")
PORT = os.getenv("PORT")
PROJECT_ROOT = os.path.dirname(os.path.abspath(__file__))
DATA_PATH = os.getenv("DATA_PATH", os.path.join(PROJECT_ROOT, "Docs"))
GURUCOOL_API_KEY = os.getenv("GURUCOOL_API_KEY")

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

app = Flask(__name__)
CORS(app, supports_credentials=True)
app.secret_key = GURUCOOL_API_KEY

cred = credentials.Certificate(CREDENTIALS_PATH)
firebase_admin.initialize_app(cred, {
    'storageBucket': STORAGE_BUCKET  
})

db = firestore.client()
storage = firebase_storage.bucket()

ollama_model = OllamaLLM(model="llama3.1")

# --------------------------------------------------------------- API Key Middleware --------------------------------------------------------------- 

def require_api_key(func):
    def wrapper(*args, **kwargs):
        api_key = request.headers.get('GuruCool-API-Key')
        if api_key != GURUCOOL_API_KEY:
            logging.error(f"Unauthorized access attempt with API Key: {api_key}")
            abort(401, description="Unauthorized access. Invalid API key.")
        return func(*args, **kwargs)
    wrapper.__name__ = func.__name__ 
    return wrapper

# --------------------------------------------------------------- EMBEDDING --------------------------------------------------------------- 
def download_files(documents, save_directory):
    """
    Downloads files from the provided document data and saves them in the specified directory.
    Args:
        documents (list): A list of dictionaries containing document names and URLs.
        save_directory (str): The directory where the files will be saved.
    """
    logging.info(f"Downloading files from provided URLs")
    start_time = time.time()

    for doc in documents:
        try:
            file_url = doc.get('url')
            file_name = doc.get('name')

            if not file_url or not file_name:
                logging.error("Invalid document data: URL or name is empty.")
                continue
            
            # Ensure file name has an extension, defaulting to .txt if not
            if '.' not in file_name:
                file_name += '.txt'
            
            file_path = os.path.join(save_directory, file_name)

            response = requests.get(file_url)
            response.raise_for_status()

            with open(file_path, 'wb') as f:
                f.write(response.content)
            logging.info(f"Downloaded and saved file: {file_name} at {file_path}")
        except requests.exceptions.RequestException as e:
            logging.error(f"Error downloading file from {file_url}: {e}")
            continue

    end_time = time.time()
    logging.info(f"Files downloaded in {end_time - start_time:.2f} seconds")

def load_pptx_files(directory):
    logging.info(f"Loading PPTX files from directory: {directory}")
    pptx_documents = []
    for pptx_file in Path(directory).glob("*.pptx"):
        logging.info(f"Processing PPTX file: {pptx_file}")
        prs = Presentation(pptx_file)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        pptx_documents.append(Document(page_content="\n".join(text_runs), metadata={"source": str(pptx_file)}))
    logging.info(f"Loaded {len(pptx_documents)} PPTX documents.")
    return pptx_documents

def load_documents(directory):
    logging.info(f"Loading documents from directory: {directory}")
    start_time = time.time()
    
    # Convert directory path to Path object
    dir_path = Path(directory)
    
    # Initialize lists to hold documents
    txt_documents = []
    pdf_documents = []
    
    # Load text files
    for file in dir_path.glob("*.txt"):
        try:
            loader = TextLoader(str(file))
            txt_documents.extend(loader.load())
            logging.info(f"Loaded text document: {file}")
        except Exception as e:
            logging.error(f"Error loading text file {file}: {e}")
    
    # Load PDF files
    pdf_loader = DirectoryLoader(str(dir_path), glob="*.pdf", loader_cls=PyPDFLoader)
    pdf_documents = pdf_loader.load()
    logging.info(f"Loaded {len(pdf_documents)} PDF documents.")
    
    # Load PPTX files
    pptx_documents = load_pptx_files(directory)
    
    # Combine all documents
    documents = txt_documents + pdf_documents + pptx_documents
    logging.info(f"Total documents loaded: {len(documents)}")
    
    # Split documents into chunks
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
    split_documents = text_splitter.split_documents(documents)
    logging.info(f"Documents split into {len(split_documents)} chunks.")
    
    end_time = time.time()
    logging.info(f"Documents loaded and processed in {end_time - start_time:.2f} seconds")

    return split_documents
def generate_and_save_embeddings(lecture_id, documents):
    cache_directory = os.path.join(PROJECT_ROOT, "Cache")
    os.makedirs(cache_directory, exist_ok=True)  

    embedding_cache_file = os.path.join(cache_directory, f"{lecture_id}_embedding_cache.pkl")
    
    logging.info(f"Checking for cached embeddings at: {embedding_cache_file}")
    if os.path.exists(embedding_cache_file):
        logging.info("Embeddings cache already exists. No need to regenerate.")
        return embedding_cache_file
    
    logging.info("No cache file found. Generating embeddings...")
    start_time = time.time()

    embedding_model = SentenceTransformer('all-MiniLM-L6-v2')
    document_texts = [doc.page_content for doc in documents]
    
    embeddings = embedding_model.encode(document_texts, convert_to_tensor=False)
    embedding_dim = len(embeddings[0])
    index = faiss.IndexFlatL2(embedding_dim)
    index.add(embeddings)
    
    try:
        with open(embedding_cache_file, 'wb') as f:
            pickle.dump((embeddings, index, documents), f)
        logging.info(f"Embeddings cached successfully with documents at {embedding_cache_file}.")
    except Exception as e:
        logging.error(f"Error saving embeddings to cache: {e}")
        raise

    end_time = time.time()
    logging.info(f"Embeddings generated and saved in {end_time - start_time:.2f} seconds")
    
    return embedding_cache_file

# --------------------------------------------------------------- TRANSCRIBER --------------------------------------------------------------- 

class Transcriber:
    def __init__(self, model_name):
        try:
            self.model = whisper.load_model(model_name)
            logging.info(f"Model '{model_name}' loaded successfully.")
        except Exception as e:
            logging.error(f"Failed to load model: {e}")
            sys.exit(1)

    def transcribe_audio(self, audio_file):
        try:
            logging.info(f"Transcribing {audio_file}...")
            result = self.model.transcribe(audio_file)
            transcript = result['text']
            if transcript:
                logging.info(f"Transcription successful: {len(transcript)} characters")
            else:
                logging.error("Transcription returned None or empty result.")
            return transcript
        except Exception as e:
            logging.error(f"Error during transcription: {e}")
            return None

def download_file_from_url(url, save_directory):
    try:
        logging.info(f"Attempting to download file from {url}")
        response = requests.get(url)
        response.raise_for_status()  
        file_name = os.path.basename(url.split('?')[0]) 
        logging.info(f"File name is - {file_name}")
        file_path = os.path.join(save_directory, file_name)
 
        with open(file_path, "wb") as file:
            file.write(response.content)
        
        logging.info(f"File downloaded successfully and saved to {file_path}")
        return file_path
    except Exception as e:
        logging.error(f"Error downloading file: {e}")
        raise

def upload_transcript_to_storage(transcript_text, firebase_storage_path):
    try:
        logging.info(f"Uploading transcript to Firebase Storage at {firebase_storage_path}")
        blob = storage.blob(firebase_storage_path)
        blob.upload_from_string(transcript_text.encode('utf-8'), content_type='text/plain')
        
        blob.make_public()
        public_url = blob.public_url
        logging.info(f"Transcript uploaded successfully. Public URL: {public_url}")
        return public_url
    except Exception as e:
        logging.error(f"Error uploading transcript to Firebase Storage: {e}")
        raise

def summarize_transcript(subject, transcript):
    prompt = f"""
    You are a helpful assistant that summarizes transcripts. Your name is HelpGuru. The transcript is from a lecture about {subject}.
    Here is the transcript:

    {transcript}

    Please provide a concise summary of this lecture with bullet points summarizing the key highlights.
    """
    
    result = ollama_model.invoke(input=prompt)
    return result

# --------------------------------------------------------------- QUIZ ---------------------------------------------------------------

def get_embedding_cache_file(lecture_id):
    return os.path.join(PROJECT_ROOT, f"Cache/{lecture_id}_embedding_cache.pkl")

def load_embeddings_and_documents(lecture_id, subject, unit, document_urls=None):
    """
    Loads embeddings and documents for a given lecture.
    Args:
        lecture_id (str): The ID of the lecture.
        subject (str): The subject of the lecture.
        unit (str): The unit of the lecture.
        document_urls (list, optional): List of document data (name and URL) for the lecture documents. Defaults to None.
    Returns:
        list: The loaded documents.
    Raises:
        ValueError: If document URLs are required but not found.
    """
    try:
        embedding_cache_file = get_embedding_cache_file(lecture_id)
        logging.info(f"Loading embeddings from cache at: {embedding_cache_file}")
        
        if not os.path.exists(embedding_cache_file):
            if not document_urls:
                logging.info(f"No document data provided. Fetching from Firestore based on Lecture ID: {lecture_id}")
                document_urls = fetch_document_urls_from_firestore(subject, unit, lecture_id)
                
                if not document_urls:
                    logging.error("Document data is required to generate embeddings but was not found.")
                    raise ValueError("Document data is required to generate embeddings but was not found.")
            
            logging.info(f"Embeddings not found for lecture ID: {lecture_id}. Generating embeddings...")
            lecture_dir = os.path.join(DATA_PATH, subject, unit, lecture_id)
            os.makedirs(lecture_dir, exist_ok=True)
            
            download_files(document_urls, lecture_dir)
            documents = load_documents(lecture_dir)
            if not documents:
                logging.error("No documents loaded, cannot generate embeddings.")
                return []
            
            embedding_cache_file = generate_and_save_embeddings(lecture_id, documents)
            logging.info(f"Embeddings generated and saved for lecture ID: {lecture_id}")
        
        with open(embedding_cache_file, 'rb') as f:
            embeddings, index, documents = pickle.load(f)
        
        logging.info("Embeddings, index, and documents loaded successfully.")
        return documents
    
    except Exception as e:
        logging.error(f"Error loading embeddings and documents: {e}")
        return []

def chunk_content(content, chunk_size=1000):
    return [content[i:i+chunk_size] for i in range(0, len(content), chunk_size)]

def generate_single_question(llm, content_chunk, question_number):
    prompt = (
        f"Using the content below, generate exactly one quiz question with an answer. "
        "Do NOT include any explanations, summaries, or extra information. "
        "The output must strictly follow this format:\n\n"
        f"Question {question_number}: <Insert question>\n"
        f"Answer {question_number}: <Insert answer>\n\n"
        "Make sure that there is nothing in the output other than the question and answer.\n\n"
        f"Content:\n{content_chunk}"
    )
    
    response = llm.generate([prompt], clean_up_tokenization_spaces=True)
    quiz_text = "".join(chunk.text for chunk in response.generations[0])
    
    return quiz_text

def generate_quiz_using_llm(documents, num_questions=10):
    try:
        if not documents:
            raise ValueError("No documents loaded.")
        
        ollama_model = OllamaLLM(model="llama3.1")
        
        combined_content = "\n\n".join(doc.page_content for doc in documents)
        content_chunks = chunk_content(combined_content)
        
        quiz = {}
        
        for i in range(min(num_questions, len(content_chunks))):
            quiz_text = generate_single_question(ollama_model, content_chunks[i], i+1)
            quiz.update(extract_questions_from_response(quiz_text))
        
        return quiz
    
    except Exception as e:
        print(f"Error generating quiz: {e}")
        return {}

def extract_questions_from_response(response_text):
    quiz = {}
    current_question = None
    current_answer = None
    
    lines = response_text.splitlines()
    
    for line in lines:
        line = line.strip()
        if line.lower().startswith("question"):
            if current_question and current_answer:
                quiz[current_question] = current_answer
            current_question = line.split(":")[1].strip()
            current_answer = None
        elif line.lower().startswith("answer"):
            current_answer = line.split(":")[1].strip()
    
    if current_question and current_answer:
        quiz[current_question] = current_answer
    
    return quiz

# --------------------------------------------------------------- ASK GURU  ---------------------------------------------------------------

def get_embedding_cache_file(lecture_id):
    """
    Returns the file path for the embedding cache file of a given lecture.

    Parameters:
    lecture_id (str): The ID of the lecture.

    Returns:
    str: The file path for the embedding cache file.
    """
    return os.path.join(PROJECT_ROOT, f"Cache/{lecture_id}_embedding_cache.pkl")

def load_embeddings(lecture_id, subject, unit, document_urls=None):
    """
    Load embeddings from cache file or generate them if not found.
    Parameters:
    lecture_id (str): The ID of the lecture.
    subject (str): The subject of the lecture.
    unit (str): The unit of the lecture.
    document_urls (list, optional): List of URLs of the lecture documents. Defaults to None.
    Returns:
    tuple: A tuple containing the embeddings, index, and documents.
    Raises:
    FileNotFoundError: If the embedding cache file is not found for the given lecture ID and no document URLs are provided.
    """
    embedding_cache_file = get_embedding_cache_file(lecture_id)
    logging.info(f"Loading embeddings from cache at: {embedding_cache_file}")
    
    if not os.path.exists(embedding_cache_file):
        if document_urls:
            logging.info(f"Embeddings not found for lecture ID: {lecture_id}. Generating embeddings...")
            

            lecture_dir = os.path.join(DATA_PATH, subject, unit, lecture_id)
            os.makedirs(lecture_dir, exist_ok=True)
            
            download_files(document_urls, lecture_dir)
            documents = load_documents(lecture_dir)
            embedding_cache_file = generate_and_save_embeddings(lecture_id, documents)
            logging.info(f"Embeddings generated and saved for lecture ID: {lecture_id}")
        else:
            raise FileNotFoundError(f"Embedding cache file not found for lecture ID: {lecture_id} and no document URLs provided.")
    
    with open(embedding_cache_file, 'rb') as f:
        embeddings, index, documents = pickle.load(f)
    
    logging.info(f"Embeddings, index, and documents loaded successfully for lecture ID: {lecture_id}.")
    return embeddings, index, documents

def fetch_document_urls_from_firestore(subject_id, unit_id, lecture_id):
    """
    Fetches the document URLs and names from Firestore for a given lecture.
    Args:
        subject_id (str): The ID of the subject.
        unit_id (str): The ID of the unit.
        lecture_id (str): The ID of the lecture.
    Returns:
        list: A list of dictionaries containing document names and URLs.
    Raises:
        Exception: If there is an error fetching the document data from Firestore.
    """
    try:
        doc_ref = db.collection('subjects').document(subject_id)\
                    .collection('units').document(unit_id)\
                    .collection('lectures').document(lecture_id)
        lecture_doc = doc_ref.get()
        
        if lecture_doc.exists:
            lecture_data = lecture_doc.to_dict()
            print(lecture_data)
            
            # Extract the relevant document URLs from `documentUrls`
            documents = lecture_data.get('documentUrls', [])
            
            if documents:
                logging.info(f"Found {len(documents)} documents in Firestore for Lecture ID: {lecture_id}")
                return documents
            else:
                logging.warning(f"No documents found in Firestore for Lecture ID: {lecture_id}")
        else:
            logging.warning(f"Lecture document not found in Firestore for Lecture ID: {lecture_id}")
    except Exception as e:
        logging.error(f"Error fetching documents from Firestore: {e}")
    
    return []



def load_embeddings_and_documents(lecture_id, subject, unit, document_urls=None):
    """
    Loads embeddings and documents for a given lecture.
    Args:
        lecture_id (str): The ID of the lecture.
        subject (str): The subject of the lecture.
        unit (str): The unit of the lecture.
        document_urls (list, optional): List of URLs for the lecture documents. Defaults to None.
    Returns:
        list: The loaded documents.
    Raises:
        ValueError: If document URLs are required but not found.
    """
    try:
        embedding_cache_file = get_embedding_cache_file(lecture_id)
        logging.info(f"Loading embeddings from cache at: {embedding_cache_file}")
        
        if not os.path.exists(embedding_cache_file):
            if not document_urls or any(not url for url in document_urls):
                logging.info(f"No document URLs provided. Fetching from Firestore based on Lecture ID: {lecture_id}")
                document_urls = fetch_document_urls_from_firestore(subject, unit, lecture_id)
                
                if not document_urls:
                    logging.error("Document URLs are required to generate embeddings but were not found.")
                    raise ValueError("Document URLs are required to generate embeddings but were not found.")
            
            logging.info(f"Embeddings not found for lecture ID: {lecture_id}. Generating embeddings...")
            lecture_dir = os.path.join(DATA_PATH, subject, unit, lecture_id)
            os.makedirs(lecture_dir, exist_ok=True)
            
            download_files(document_urls, lecture_dir)
            documents = load_documents(lecture_dir)
            if not documents:
                logging.error("No documents loaded, cannot generate embeddings.")
                return []
            
            embedding_cache_file = generate_and_save_embeddings(lecture_id, documents)
            logging.info(f"Embeddings generated and saved for lecture ID: {lecture_id}")
        
        with open(embedding_cache_file, 'rb') as f:
            embeddings, index, documents = pickle.load(f)
        
        logging.info("Embeddings, index, and documents loaded successfully.")
        return documents
    
    except Exception as e:
        logging.error(f"Error loading embeddings and documents: {e}")
        return []


def retrieve_relevant_docs(query, index, documents, embedding_model):
    """
    Retrieves relevant documents based on a given query.

    Args:
        query (str): The query string.
        index (Index): The index used for searching.
        documents (list): The list of documents.
        embedding_model (EmbeddingModel): The embedding model used for encoding the query.

    Returns:
        list: The list of relevant documents.
    """
    query_embedding = embedding_model.encode([query], convert_to_tensor=False)
    _, indices = index.search(query_embedding, k=5)
    relevant_docs = [documents[i] for i in indices[0]] 
    return relevant_docs

def generate_response(query, relevant_docs, conversation_history):
    """
    Generates a response to a given query based on the provided context and relevant documents.
    Args:
        query (str): The question/query to be answered.
        relevant_docs (list): A list of relevant documents.
        conversation_history (list): A list of tuples representing the conversation history.
    Returns:
        tuple: A tuple containing the generated answer and a list of document sources.
    Raises:
        Exception: If there is an error generating the response.
    """
    context = " ".join([f"Q: {q}\nA: {a}" for q, a in conversation_history])
    context += " " + " ".join([doc.page_content for doc in relevant_docs])
    
    prompt = (
        f"Answer the following question based on the provided context. If the answer is not explicitly available in the context, "
        f"you may still provide a general response, but specify that the information could not be found in the provided documents.\n\n"
        f"Question: {query}\n"
        f"Context: {context}\n"
        f"Answer:"
    )
    
    try:
        response = ollama_model.generate([prompt], clean_up_tokenization_spaces=True)
        answer = "".join(chunk.text for chunk in response.generations[0])
        logging.info(f"Ollama response generated successfully.")
    except Exception as e:
        logging.error(f"Error generating response: {e}")
        answer = "Sorry, I couldn't generate a response."
    
    document_sources = [doc.metadata.get("source") for doc in relevant_docs]
    return answer, document_sources
#  --------------------------------------------------------------- ROUTES --------------------------------------------------------------- 

@app.route('/', methods=['GET'])
def hello_world():
    return '''
    <html>
        <head>
            <title>Hello World</title>
        </head>
        <body>
            <h1>Hello, World!</h1>
            <p>GuruCool API is functional</p>
            <p>API Key is: 1234</p>
            <p>lmao</p>
        </body>
    </html>
    '''

@app.route('/generate_embeddings', methods=['POST'])
@require_api_key
def generate_embeddings():
    try:
        data = request.json
        subject = data.get('subject')
        unit = data.get('unit')
        lecture_id = data.get('lecture')
        document_urls = data.get('document_urls', [])
        print(document_urls)

        if not subject or not unit or not lecture_id or not document_urls:
            logging.error("Missing fields: 'subject', 'unit', 'lecture', and 'document_urls' are required.")
            return jsonify({"error": "Missing fields: 'subject', 'unit', 'lecture', and 'document_urls' are required."}), 400
        
        logging.info(f"Starting embedding generation for lecture ID: {lecture_id}")

        lecture_dir = os.path.join(DATA_PATH, subject, unit, lecture_id)
        os.makedirs(lecture_dir, exist_ok=True)
        
        download_files(document_urls, lecture_dir)
        documents = load_documents(lecture_dir)
        print("The documents fetched are -",documents)
        embedding_cache_file = generate_and_save_embeddings(lecture_id, documents)

        logging.info(f"Embedding generation completed for lecture ID: {lecture_id}")
        return jsonify({"message": "Embeddings generated and cached successfully", "embedding_cache_file": embedding_cache_file}), 200

    except Exception as e:
        logging.error(f"An error occurred: {e}")
        return jsonify({"error": str(e)}), 500

    
@app.route('/transcribe_audio', methods=['POST'])
@require_api_key
def transcribe_audio():
    logging.info(f"Request received with body: {request.json}")
    try:
        data = request.json
        file_url = data.get('download_url')
        model_name = data.get('model_name', 'base')
        subject = data.get('subject')
        unit = data.get('unit')
        lecture = data.get('lecture')

        base_dir = os.path.join("Transcript", subject, unit, lecture)
        audio_dir = os.path.join(base_dir, "audio")
        transcript_dir = os.path.join(base_dir, "transcripts")
        
        logging.info(f"Attempting to download file from {file_url}")

        os.makedirs(audio_dir, exist_ok=True)
        os.makedirs(transcript_dir, exist_ok=True)

        audio_file_path = download_file_from_url(file_url, audio_dir)
        
        transcriber = Transcriber(model_name=model_name)
        transcript = transcriber.transcribe_audio(audio_file_path)
        
        if transcript:
            logging.info(f"Transcript generated successfully.")
            logging.info(f"Transcript (first 500 chars): {transcript[:500]}")

            transcript_file_name = os.path.splitext(os.path.basename(audio_file_path))[0] + ".txt"
            local_transcript_path = os.path.join(transcript_dir, transcript_file_name)
            with open(local_transcript_path, "w", encoding="utf-8") as transcript_file:
                transcript_file.write(transcript)
            logging.info(f"Transcript saved locally at {local_transcript_path}")

            summary = summarize_transcript(subject, transcript)
            summary_file_name = "summary_" + transcript_file_name
            local_summary_path = os.path.join(transcript_dir, summary_file_name)
            with open(local_summary_path, "w", encoding="utf-8") as summary_file:
                summary_file.write(summary)
            logging.info(f"Summary saved locally at {local_summary_path}")

            firebase_storage_path = f"{subject}/{unit}/{lecture}/{transcript_file_name}"
            transcript_url = upload_transcript_to_storage(transcript, firebase_storage_path)

            summary_storage_path = f"{subject}/{unit}/{lecture}/{summary_file_name}"
            summary_url = upload_transcript_to_storage(summary, summary_storage_path)
            logging.info(f"Transcript uploaded. Public URL: {transcript_url}")
            logging.info(f"Summary uploaded. Public URL: {summary_url}")

            doc_ref = db.collection('subjects').document(subject)\
                        .collection('units').document(unit)\
                        .collection('lectures').document(lecture)

            doc_ref.update({
                'transcript_url': transcript_url,
                'summary_url': summary_url,
                'documentUrls': firestore.ArrayUnion([
                    {"name": "Transcript", "url": transcript_url},
                    {"name": "Summary", "url": summary_url}
                ])
            })
            
            logging.info(f"Firestore document updated successfully for lecture {lecture}")
            
            return jsonify({"message": "Transcription and summary successful", "transcript_url": transcript_url, "summary_url": summary_url}), 200
        else:
            logging.error("Transcription failed.")
            return jsonify({"error": "Transcription failed."}), 500
    
    except Exception as e:
        logging.error(f"An error occurred: {e}")
        return jsonify({"error": str(e)}), 500
    
@app.route('/generate_quiz', methods=['POST'])
@require_api_key
def generate_quiz():
    try:
        data = request.json
        subject = data.get("subject_id")
        unit = data.get("unit_id")
        lecture_id = data.get("lecture_id")
        document_urls = data.get("document_urls", [])
        num_questions = int(data.get("num_questions", 10))
        
        if not all([subject, unit, lecture_id]):
            logging.error("Required fields 'subject', 'unit', and 'lecture_id' are missing.")
            return jsonify({"error": "Required fields 'subject', 'unit', and 'lecture_id' are missing."}), 400

        logging.info(f"Generating quiz for Lecture ID: {lecture_id}, Subject: {subject}, Unit: {unit}")

        print(lecture_id, subject, unit, document_urls)
        documents = load_embeddings_and_documents(lecture_id, subject, unit, document_urls)
        
        if not documents:
            logging.error(f"No documents found for Lecture ID: {lecture_id}, Subject: {subject}, Unit: {unit}")
            return jsonify({"error": "No documents found for the given lecture_id"}), 404
        
        quiz = generate_quiz_using_llm(documents, num_questions)
        
        if quiz:
            logging.info(f"Quiz generated successfully for Lecture ID: {lecture_id}")
            return jsonify({"quiz": quiz}), 200
        else:
            logging.error("No quiz generated. Please check the input content or LLM output.")
            return jsonify({"error": "No quiz generated. Please check the input content or LLM output."}), 500
    
    except ValueError as e:
        logging.error(f"ValueError: {e}")
        return jsonify({"error": str(e)}), 400
    except Exception as e:
        logging.error(f"Internal error: {e}")
        return jsonify({"error": "An internal error occurred"}), 500

    
@app.route('/chat', methods=['POST'])
@require_api_key
def chat():
    start_time = time.time()
    
    if 'conversation_history' not in session:
        session['conversation_history'] = []
    
    data = request.json
    user_query = request.json.get("query")
    lecture_id = request.json.get("lecture_id")
    subject = data.get("subject_id")
    unit = data.get("unit_id")
    
    if not user_query or not lecture_id:
        logging.error("Both 'query' and 'lecture_id' are required.")
        return jsonify({"error": "Both 'query' and 'lecture_id' are required."}), 400
    
    logging.info(f"Received query: {user_query}, Lecture ID: {lecture_id}")
    
    try:
        embeddings, index, documents = load_embeddings(lecture_id, subject, unit)
    except FileNotFoundError as e:
        logging.error(f"Error loading embeddings: {e}")
        return jsonify({"error": str(e)}), 500
    
    embedding_model = SentenceTransformer('all-MiniLM-L6-v2')

    retrieval_start_time = time.time()
    relevant_docs = retrieve_relevant_docs(user_query, index, documents, embedding_model)
    retrieval_end_time = time.time()
    
    response, document_sources = generate_response(user_query, relevant_docs, session['conversation_history'])

    session['conversation_history'].append((user_query, response))
    
    end_time = time.time()

    logging.info(f"Document retrieval time: {retrieval_end_time - retrieval_start_time:.2f} seconds")
    logging.info(f"Total processing time: {end_time - start_time:.2f} seconds")
    logging.info(f"Referenced Documents: {', '.join(document_sources)}")
    
    return jsonify({
        "response": response,
        "document_sources": document_sources,
        "processing_time": f"{end_time - start_time:.2f} seconds"
    })

@app.route('/reset', methods=['POST'])
def reset():
    session.pop('conversation_history', None)
    logging.info("Conversation history reset.")
    return jsonify({"message": "Conversation history reset."})

# --------------------------------------------------------------- MAIN --------------------------------------------------------------- 

if __name__ == "__main__":  
    app.run(debug=True, host='0.0.0.0', port=PORT)
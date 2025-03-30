import os
import pickle
import requests
import time
from flask import Flask, request, jsonify
from dotenv import load_dotenv
from langchain_community.document_loaders import DirectoryLoader, PyPDFLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document
from sentence_transformers import SentenceTransformer
import faiss
from pptx import Presentation
from pathlib import Path
import warnings
import firebase_admin
from firebase_admin import credentials, firestore
import logging
from flask_cors import CORS

warnings.filterwarnings("ignore", category=FutureWarning)

load_dotenv()
PROJECT_ROOT = os.path.dirname(os.path.abspath(__file__))
DATA_PATH = os.getenv("DATA_PATH", os.path.join(PROJECT_ROOT, "Docs"))
STORAGE_BUCKET = os.getenv("STORAGE_BUCKET")
CREDENTIALS_PATH = os.getenv("CREDENTIALS_PATH")

app = Flask(__name__)
CORS(app)

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
cred = credentials.Certificate(CREDENTIALS_PATH)
firebase_admin.initialize_app(cred, {'storageBucket': STORAGE_BUCKET})
db = firestore.client()

def download_files(document_urls, save_directory):
    logging.info(f"Downloading files from provided URLs")
    start_time = time.time()

    for file_url in document_urls:
        try:
            if not file_url:
                logging.error("Invalid URL provided: URL is empty.")
                continue
            
            file_name = os.path.basename(file_url.split('?')[0])
            file_path = os.path.join(save_directory, file_name)

            response = requests.get(file_url)
            response.raise_for_status()  # Check for HTTP errors

            with open(file_path, 'wb') as f:
                f.write(response.content)
            logging.info(f"Downloaded and saved file: {file_name} at {file_path}")
        except requests.exceptions.RequestException as e:
            logging.error(f"Error downloading file from {file_url}: {e}")
            continue

    end_time = time.time()
    logging.info(f"Files downloaded in {end_time - start_time:.2f} seconds")

def load_pptx_files(directory):
    logging.info(f"Loading PPTX files from directory: {directory}")
    pptx_documents = []
    for pptx_file in Path(directory).glob("*.pptx"):
        logging.info(f"Processing PPTX file: {pptx_file}")
        prs = Presentation(pptx_file)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        pptx_documents.append(Document(page_content="\n".join(text_runs), metadata={"source": str(pptx_file)}))
    logging.info(f"Loaded {len(pptx_documents)} PPTX documents.")
    return pptx_documents

def load_documents(directory):
    logging.info(f"Loading documents from directory: {directory}")
    start_time = time.time()

    txt_loader = DirectoryLoader(directory, glob="*.txt")
    txt_documents = txt_loader.load()
    logging.info(f"Loaded {len(txt_documents)} text documents.")
    
    pdf_loader = DirectoryLoader(directory, glob="*.pdf", loader_cls=PyPDFLoader)
    pdf_documents = pdf_loader.load()
    logging.info(f"Loaded {len(pdf_documents)} PDF documents.")
    
    pptx_documents = load_pptx_files(directory)
    
    documents = txt_documents + pdf_documents + pptx_documents
    logging.info(f"Total documents loaded: {len(documents)}")
    
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
    split_documents = text_splitter.split_documents(documents)
    logging.info(f"Documents split into {len(split_documents)} chunks.")
    
    end_time = time.time()
    logging.info(f"Documents loaded and processed in {end_time - start_time:.2f} seconds")

    return split_documents

def generate_and_save_embeddings(lecture_id, documents):
    cache_directory = os.path.join(PROJECT_ROOT, "Cache")
    os.makedirs(cache_directory, exist_ok=True)  

    embedding_cache_file = os.path.join(cache_directory, f"{lecture_id}_embedding_cache.pkl")
    
    logging.info(f"Checking for cached embeddings at: {embedding_cache_file}")
    if os.path.exists(embedding_cache_file):
        logging.info("Embeddings cache already exists. No need to regenerate.")
        return embedding_cache_file
    
    logging.info("No cache file found. Generating embeddings...")
    start_time = time.time()

    embedding_model = SentenceTransformer('all-MiniLM-L6-v2')
    document_texts = [doc.page_content for doc in documents]
    
    embeddings = embedding_model.encode(document_texts, convert_to_tensor=False)
    embedding_dim = len(embeddings[0])
    index = faiss.IndexFlatL2(embedding_dim)
    index.add(embeddings)
    
    try:
        with open(embedding_cache_file, 'wb') as f:
            pickle.dump((embeddings, index, documents), f)
        logging.info(f"Embeddings cached successfully with documents at {embedding_cache_file}.")
    except Exception as e:
        logging.error(f"Error saving embeddings to cache: {e}")
        raise

    end_time = time.time()
    logging.info(f"Embeddings generated and saved in {end_time - start_time:.2f} seconds")
    
    return embedding_cache_file


@app.route('/generate_embeddings', methods=['POST'])
def generate_embeddings():
    try:
        subject = request.headers.get('subject')
        unit = request.headers.get('unit')
        lecture_id = request.headers.get('lecture')
        document_urls = request.headers.get('document_urls', '').split(',')

        if not subject or not unit or not lecture_id or not document_urls:
            logging.error("Missing headers: 'subject', 'unit', 'lecture', and 'document_urls' are required.")
            return jsonify({"error": "Missing headers: 'subject', 'unit', 'lecture', and 'document_urls' are required."}), 400
        
        logging.info(f"Starting embedding generation for lecture ID: {lecture_id}")

        lecture_dir = os.path.join(DATA_PATH, subject, unit, lecture_id)
        os.makedirs(lecture_dir, exist_ok=True)
        
        download_files(document_urls, lecture_dir)
        documents = load_documents(lecture_dir)
        embedding_cache_file = generate_and_save_embeddings(lecture_id, documents)

        logging.info(f"Embedding generation completed for lecture ID: {lecture_id}")
        return jsonify({"message": "Embeddings generated and cached successfully", "embedding_cache_file": embedding_cache_file}), 200

    except Exception as e:
        logging.error(f"An error occurred: {e}")
        return jsonify({"error": str(e)}), 500

if __name__ == "__main__":
    app.run(debug=True, host='0.0.0.0', port=8000)

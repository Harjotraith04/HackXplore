import os
import pickle
from dotenv import load_dotenv
from langchain_community.document_loaders import DirectoryLoader, PyPDFLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document
from langchain_ollama import OllamaLLM
from sentence_transformers import SentenceTransformer
import faiss
from pptx import Presentation
from pathlib import Path
import time
import warnings
warnings.filterwarnings("ignore", category=FutureWarning)

load_dotenv()

DATA_PATH = os.getenv("DATA_PATH", "Docs/")
EMBEDDING_CACHE_FILE = os.getenv("EMBEDDING_CACHE_FILE", "embedding_cache.pkl")

def load_pptx_files(directory):
    pptx_documents = []
    for pptx_file in Path(directory).glob("*.pptx"):
        prs = Presentation(pptx_file)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        pptx_documents.append(Document(page_content="\n".join(text_runs), metadata={"source": str(pptx_file)}))
    return pptx_documents

def load_documents():
    start_time = time.time()
    print("Loading documents...")
    
    txt_loader = DirectoryLoader(DATA_PATH, glob="*.txt")
    txt_documents = txt_loader.load()
    
    pdf_loader = DirectoryLoader(DATA_PATH, glob="*.pdf", loader_cls=PyPDFLoader)
    pdf_documents = pdf_loader.load()
    
    pptx_documents = load_pptx_files(DATA_PATH)
    
    documents = txt_documents + pdf_documents + pptx_documents
    
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
    split_documents = text_splitter.split_documents(documents)
    
    print(f"Documents loaded and split into chunks. Time taken: {time.time() - start_time:.2f} seconds.")
    return split_documents

def load_or_generate_embeddings(documents):
    start_time = time.time()
    
    if os.path.exists(EMBEDDING_CACHE_FILE):
        print("Loading cached embeddings...")
        with open(EMBEDDING_CACHE_FILE, 'rb') as f:
            embeddings, index = pickle.load(f)
        print(f"Cached embeddings loaded. Time taken: {time.time() - start_time:.2f} seconds.")
    else:
        print("Generating embeddings...")
        embedding_model = SentenceTransformer('all-MiniLM-L6-v2')
        document_texts = [doc.page_content for doc in documents]
        embeddings = embedding_model.encode(document_texts, convert_to_tensor=False)
        
        embedding_dim = len(embeddings[0])
        index = faiss.IndexFlatL2(embedding_dim)
        index.add(embeddings)
        
        with open(EMBEDDING_CACHE_FILE, 'wb') as f:
            pickle.dump((embeddings, index), f)
        
        print(f"Embeddings generated and cached. Time taken: {time.time() - start_time:.2f} seconds.")
    
    return embeddings, index

def retrieve_relevant_docs(query, index, documents, embedding_model):
    start_time = time.time()
    print("Retrieving relevant documents...")
    
    query_embedding = embedding_model.encode([query], convert_to_tensor=False)
    _, indices = index.search(query_embedding, k=5)
    relevant_docs = [documents[i] for i in indices[0]]
    
    print(f"Relevant documents retrieved. Time taken: {time.time() - start_time:.2f} seconds.")
    return relevant_docs

def generate_response(query, relevant_docs, conversation_history):
    start_time = time.time()
    print("Generating response from the chatbot...")
    
    context = " ".join([f"Q: {q}\nA: {a}" for q, a in conversation_history])
    context += " " + " ".join([doc.page_content for doc in relevant_docs])
    
    prompt = f"Question: {query}\nContext: {context}\nAnswer:"
    
    ollama_model = OllamaLLM(model="llama3.1")
    response = ollama_model.generate([prompt], clean_up_tokenization_spaces=True)
    
    try:
        answer = "".join(chunk.text for chunk in response.generations[0])
        print(f"Response generated. Time taken: {time.time() - start_time:.2f} seconds.")
    except (AttributeError, IndexError):
        answer = "Sorry, I couldn't generate a response."
        print("Error in generating response.")
    
    document_sources = [doc.metadata.get("source") for doc in relevant_docs]
    return answer, document_sources

def chatbot():
    documents = load_documents()
    embedding_model = SentenceTransformer('all-MiniLM-L6-v2')
    embeddings, index = load_or_generate_embeddings(documents)
    
    conversation_history = []
    
    while True:
        query = input("You: ")
        if query.lower() in ["exit", "quit"]:
            break
        
        relevant_docs = retrieve_relevant_docs(query, index, documents, embedding_model)
        response, document_sources = generate_response(query, relevant_docs, conversation_history)
        
        print("Thinking Ask Guru:", response)
        print("Referenced Documents:", ", ".join(document_sources))
        
        conversation_history.append((query, response))

chatbot()
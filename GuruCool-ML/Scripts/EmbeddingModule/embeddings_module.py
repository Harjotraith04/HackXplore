import os
import pickle
from dotenv import load_dotenv
from langchain_community.document_loaders import DirectoryLoader, PyPDFLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document
from sentence_transformers import SentenceTransformer
import faiss
from pptx import Presentation
from pathlib import Path
import warnings
warnings.filterwarnings("ignore", category=FutureWarning)

PROJECT_ROOT = os.path.dirname(os.path.dirname(__file__))
EMBEDDING_CACHE_FILE = os.path.join(PROJECT_ROOT, "embedding_cache.pkl")

DATA_PATH = os.getenv("DATA_PATH", "Docs/")

def load_pptx_files(directory):
    print(f"Loading PPTX files from directory: {directory}")
    pptx_documents = []
    for pptx_file in Path(directory).glob("*.pptx"):
        print(f"Processing PPTX file: {pptx_file}")
        prs = Presentation(pptx_file)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        pptx_documents.append(Document(page_content="\n".join(text_runs), metadata={"source": str(pptx_file)}))
    print(f"Loaded {len(pptx_documents)} PPTX documents.")
    return pptx_documents

def load_documents():
    print(f"Loading documents from directory: {DATA_PATH}")
    
    txt_loader = DirectoryLoader(DATA_PATH, glob="*.txt")
    txt_documents = txt_loader.load()
    print(f"Loaded {len(txt_documents)} text documents.")
    
    pdf_loader = DirectoryLoader(DATA_PATH, glob="*.pdf", loader_cls=PyPDFLoader)
    pdf_documents = pdf_loader.load()
    print(f"Loaded {len(pdf_documents)} PDF documents.")
    
    pptx_documents = load_pptx_files(DATA_PATH)
    
    documents = txt_documents + pdf_documents + pptx_documents
    print(f"Total documents loaded: {len(documents)}")
    
    # Split documents into chunks
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
    split_documents = text_splitter.split_documents(documents)
    print(f"Documents split into {len(split_documents)} chunks.")
    
    return split_documents

def generate_and_save_embeddings():
    documents = load_documents()
    
    print(f"Checking for cached embeddings at: {EMBEDDING_CACHE_FILE}")
    if os.path.exists(EMBEDDING_CACHE_FILE):
        print("Embeddings cache already exists. No need to regenerate.")
        return
    
    print("No cache file found. Generating embeddings...")
    embedding_model = SentenceTransformer('all-MiniLM-L6-v2')
    document_texts = [doc.page_content for doc in documents]
    
    embeddings = embedding_model.encode(document_texts, convert_to_tensor=False)
    embedding_dim = len(embeddings[0])
    index = faiss.IndexFlatL2(embedding_dim)
    index.add(embeddings)
    
    try:
        with open(EMBEDDING_CACHE_FILE, 'wb') as f:
            pickle.dump((embeddings, index, documents), f)
        print("Embeddings cached successfully with documents.")
    except Exception as e:
        print(f"Error saving embeddings to cache: {e}")

if __name__ == "__main__":
    generate_and_save_embeddings()